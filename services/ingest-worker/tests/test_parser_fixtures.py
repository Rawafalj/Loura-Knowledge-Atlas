from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation

from loura_ingest_worker.parsing import DoclingParser
from loura_ingest_worker.segmentation import build_segments


def _minimal_pdf(text: str) -> bytes:
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, item in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode() + item + b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode())
    output.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    return bytes(output)


def _write_fixture(path: Path, kind: str) -> None:
    if kind == "md":
        path.write_text("# Atlas fixture\n\nRetries require idempotent actions.", encoding="utf-8")
    elif kind == "html":
        path.write_text("<h1>Atlas fixture</h1><p>Retries require idempotent actions.</p>")
    elif kind == "docx":
        document = Document()
        document.add_heading("Atlas fixture", 1)
        document.add_paragraph("Retries require idempotent actions.")
        document.save(path)
    elif kind == "pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Atlas fixture"
        slide.placeholders[1].text = "Retries require idempotent actions."
        presentation.save(path)
    else:
        path.write_bytes(_minimal_pdf("Atlas fixture"))


@pytest.mark.parametrize("kind", ["md", "html", "docx", "pptx", "pdf"])
def test_docling_parser_preserves_authored_fixture_text(tmp_path: Path, kind: str) -> None:
    path = tmp_path / f"fixture.{kind}"
    _write_fixture(path, kind)

    parsed = DoclingParser().parse(path)
    segments = build_segments(parsed.blocks)

    assert parsed.parser_name == "docling"
    assert segments
    assert any("Atlas fixture" in segment.text for segment in segments)
